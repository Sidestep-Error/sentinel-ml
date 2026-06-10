"""Genererar sentinel-ml Demo Day-presentationen som en .pptx-fil.

Usage:
  python scripts/generate_presentation.py
  python scripts/generate_presentation.py --out docs/sentinel-ml-presentation.pptx
"""

from __future__ import annotations

from pathlib import Path

import typer
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Färgpalett ────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1a, 0x1a, 0x2e)   # mörkblå bakgrund
ACCENT     = RGBColor(0x16, 0x21, 0x3e)   # lite ljusare blå
TEAL       = RGBColor(0x0f, 0x3d, 0x57)   # teal för rubriker
GREEN      = RGBColor(0x4c, 0xaf, 0x50)   # grön för "klart"
ORANGE     = RGBColor(0xff, 0x98, 0x00)   # orange för varning
RED        = RGBColor(0xef, 0x53, 0x50)   # röd för attack
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xb0, 0xbe, 0xc5)
YELLOW     = RGBColor(0xff, 0xd5, 0x4f)

W = Inches(10)
H = Inches(7.5)


def _bg(slide, color: RGBColor = DARK_BG) -> None:
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height, text: str, font_size: int = 18,
         bold: bool = False, color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
         bg: RGBColor | None = None, italic: bool = False) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK_BG)

    # Tealstripe
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.8), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    _box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(1.2),
         "sentinel-ml", font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _box(slide, Inches(0.6), Inches(1.6), Inches(8.8), Inches(0.8),
         "AI-driven hotklassificering, IOC-extraktion & logganomalidetektion",
         font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    _box(slide, Inches(0.6), Inches(3.1), Inches(8.8), Inches(0.5),
         "Chas Academy · Nätverks-, OT- & AI-säkerhet · Juni 2026",
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Tech-stack badges
    badges = ["scikit-learn", "FastAPI", "spaCy", "Ollama", "Kubernetes"]
    for i, badge in enumerate(badges):
        x = Inches(1.2 + i * 1.55)
        shape = slide.shapes.add_shape(1, x, Inches(4.2), Inches(1.35), Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = badge
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE

    _box(slide, Inches(0.6), Inches(6.6), Inches(8.8), Inches(0.4),
         "Demo Day — 16 juni 2026 · 12 minuter",
         font_size=13, color=YELLOW, align=PP_ALIGN.CENTER, italic=True)


def _problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
         "Problemet", font_size=32, bold=True, color=TEAL)

    problems = [
        ("ClamAV fångar känd malware", "men ny/okänd malware passerar obemärkt"),
        ("Inga IOCs extraheras automatiskt", "IP, hash, domän och CVE måste sökas manuellt"),
        ("Serverloggar analyseras inte i realtid", "attackmönster syns inte förrän det är för sent"),
        ("Ingen kategorisering av hotrapporter", "säkerhetsanalytiker måste läsa varje rapport själv"),
    ]

    for i, (prob, sol) in enumerate(problems):
        y = Inches(1.2 + i * 1.35)
        # Röd ruta med problem
        shape = slide.shapes.add_shape(1, Inches(0.5), y, Inches(0.06), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RED
        shape.line.fill.background()

        _box(slide, Inches(0.8), y, Inches(8.5), Inches(0.45),
             prob, font_size=17, bold=True, color=WHITE)
        _box(slide, Inches(0.8), Inches(y.inches + 0.45), Inches(8.5), Inches(0.5),
             f"→ {sol}", font_size=14, color=LIGHT_GRAY)

    _box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5),
         "sentinel-ml löser detta med tre ML-spår integrerade direkt i upload-flödet.",
         font_size=15, color=YELLOW, italic=True)


def _architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
         "Arkitektur", font_size=32, bold=True, color=TEAL)

    # sentinel-upload-api box
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(3.5), Inches(1.1))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = TEAL
    _box(slide, Inches(0.55), Inches(1.25), Inches(3.4), Inches(0.5),
         "sentinel-upload-api", font_size=15, bold=True, color=WHITE)
    _box(slide, Inches(0.55), Inches(1.65), Inches(3.4), Inches(0.5),
         "FastAPI · MongoDB · ClamAV", font_size=11, color=LIGHT_GRAY)

    # Arrow
    _box(slide, Inches(4.1), Inches(1.6), Inches(1.5), Inches(0.5),
         "HTTP  →\n500 ms timeout", font_size=11, color=YELLOW, align=PP_ALIGN.CENTER)

    # sentinel-ml box
    shape2 = slide.shapes.add_shape(1, Inches(5.7), Inches(1.2), Inches(3.8), Inches(1.1))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = TEAL
    shape2.line.color.rgb = GREEN
    _box(slide, Inches(5.75), Inches(1.25), Inches(3.7), Inches(0.5),
         "sentinel-ml (port 8100)", font_size=15, bold=True, color=WHITE)
    _box(slide, Inches(5.75), Inches(1.65), Inches(3.7), Inches(0.5),
         "FastAPI · scikit-learn · spaCy · Ollama", font_size=11, color=LIGHT_GRAY)

    # Endpoints
    endpoints = [
        ("POST /predict/threat", "Hotklassificering + IOC-extraktion", GREEN),
        ("POST /predict/log-anomaly", "TF-IDF + IsolationForest", GREEN),
        ("POST /predict/liveflow", "Aggregator — allt i ett svar", YELLOW),
        ("POST /predict/cve-relevance", "SBOM-komponentmatchning", LIGHT_GRAY),
        ("POST /predict/upload-ingest", "Filmetadata → riskbedömning", LIGHT_GRAY),
    ]
    for i, (ep, desc, color) in enumerate(endpoints):
        y = Inches(2.7 + i * 0.82)
        _box(slide, Inches(0.5), y, Inches(3.8), Inches(0.4),
             ep, font_size=12, bold=True, color=color,
             bg=ACCENT)
        _box(slide, Inches(4.4), y, Inches(5.2), Inches(0.4),
             desc, font_size=12, color=LIGHT_GRAY)

    _box(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.35),
         "Kubernetes NetworkPolicy öppen · Hetzner k3s · Graceful degradation om sentinel-ml är nere",
         font_size=11, color=LIGHT_GRAY, italic=True)


def _results_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
         "ML-resultat", font_size=32, bold=True, color=TEAL)

    # Tabell Alt 6
    _box(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.45),
         "Alt 6 — Threat classifier (TF-IDF + Logistic Regression · 1 582 CTI-dokument)",
         font_size=14, bold=True, color=GREEN)

    headers = ["Klass", "Precision", "Recall", "F1", "n"]
    rows = [
        ["ransomware", "0.98", "0.98", "0.98", "47"],
        ["phishing",   "0.97", "0.97", "0.97", "76"],
        ["malware",    "0.93", "0.97", "0.95", "139"],
        ["intrusion",  "0.89", "0.85", "0.87", "48"],
        ["ddos ⚠",    "1.00", "0.43", "0.60", "7"],
        ["MACRO",      "0.955","0.841","0.875","317"],
    ]
    col_w = [Inches(2), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1)]
    col_x = [Inches(0.5), Inches(2.6), Inches(3.95), Inches(5.3), Inches(6.65)]

    for j, h in enumerate(headers):
        _box(slide, col_x[j], Inches(1.65), col_w[j], Inches(0.38),
             h, font_size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        y = Inches(2.08 + i * 0.37)
        bg = ACCENT if i % 2 == 0 else DARK_BG
        highlight = (row[0] == "MACRO")
        for j, cell in enumerate(row):
            clr = YELLOW if highlight else WHITE
            if not highlight and j in (3,) and float(cell) < 0.7:
                clr = ORANGE
            _box(slide, col_x[j], y, col_w[j], Inches(0.35),
                 cell, font_size=11, bold=highlight, color=clr,
                 align=PP_ALIGN.CENTER, bg=bg)

    # Alt 4 + Alt 3
    _box(slide, Inches(0.5), Inches(4.4), Inches(4.3), Inches(0.4),
         "Alt 4 — Log anomaly (TF-IDF + IsolationForest)", font_size=13, bold=True, color=TEAL)
    _box(slide, Inches(0.5), Inches(4.85), Inches(4.3), Inches(0.35),
         "F1: 0.370  ·  syntetisk loggdata  ·  2 metoder jämförda", font_size=12, color=LIGHT_GRAY)

    _box(slide, Inches(5.3), Inches(4.4), Inches(4.3), Inches(0.4),
         "Alt 3 — Malware metadata (Random Forest)", font_size=13, bold=True, color=TEAL)
    _box(slide, Inches(5.3), Inches(4.85), Inches(4.3), Inches(0.35),
         "F1: 0.422  ·  4× bättre än slumpen  ·  metadata-begränsning", font_size=12, color=LIGHT_GRAY)

    _box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.4),
         "DDoS: 36 träningsexempel av 1 582 — datavolymproblem, inte modellproblem.",
         font_size=12, color=ORANGE, italic=True)

    # Definitioner
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(6.05), Inches(9), Inches(0.04))
    sep.fill.solid(); sep.fill.fore_color.rgb = TEAL; sep.line.fill.background()

    defs = [
        ("Precision:", "Av flaggade hot — hur många var verkliga? (undviker falsklarm)"),
        ("Recall:", "Av alla verkliga hot — hur många hittades? (undviker att missa hot)"),
        ("F1:", "Harmoniskt medel av Precision och Recall — balanserat nyckeltal"),
        ("n:", "Antal testexempel för klassen (20 % holdout-split)"),
    ]
    for i, (term, desc) in enumerate(defs):
        y = Inches(6.15 + i * 0.31)
        _box(slide, Inches(0.5), y, Inches(1.1), Inches(0.28),
             term, font_size=10, bold=True, color=YELLOW)
        _box(slide, Inches(1.65), y, Inches(7.9), Inches(0.28),
             desc, font_size=10, color=LIGHT_GRAY)


def _adversarial_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
         "Adversarial-analys — vi attackerade vår egen modell", font_size=28, bold=True, color=RED)

    experiments = [
        ("🔴 Data Poisoning", "Spår A — Threat classifier",
         "Felmärkte 5–20 % av träningsdatat (label-flipping).\nΔF1 vid 20 %: −0.156  (0.963 → 0.807)\nGradvis degradering — syns inte utan löpande monitorering.",
         RED),
        ("🟡 Evasion", "Upload-classifier — Random Forest",
         "Uniform brusinjicering ε ∈ {0.01–0.2} på feature-vektorn.\nFlip-rate: 0.0 % — RF är robust mot slumpmässigt brus.\nSlutsats: riktad ART-attack (HopSkipJump) krävs.",
         ORANGE),
        ("🔵 Mimicry (Best Heist)", "Log-anomali — TF-IDF + IsolationForest",
         "77 % av attackloggar undkommer redan utan modifiering.\nReverse shell och path traversal är osynliga för modellen.\n'En angripare väljer rätt attacktyp — inget mimicry behövs.'",
         TEAL),
    ]

    for i, (title, subtitle, body, color) in enumerate(experiments):
        y = Inches(1.2 + i * 1.95)
        shape = slide.shapes.add_shape(1, Inches(0.4), y, Inches(9.2), Inches(1.8))
        shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT
        shape.line.color.rgb = color

        _box(slide, Inches(0.6), y + Inches(0.1), Inches(5), Inches(0.45),
             title, font_size=16, bold=True, color=color)
        _box(slide, Inches(5.8), y + Inches(0.1), Inches(3.6), Inches(0.45),
             subtitle, font_size=12, color=LIGHT_GRAY, italic=True)
        _box(slide, Inches(0.6), y + Inches(0.55), Inches(8.8), Inches(1.1),
             body, font_size=12, color=WHITE)


def _limitations_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
         "Begränsningar & nästa steg", font_size=32, bold=True, color=TEAL)

    items = [
        ("Upload-klassificerare", "Tränad på syntetisk data. Förklädda filer (exe→pdf) passerar.",
         "Shannon-entropiberäkning + magic byte-detektering på filinnehållet"),
        ("Log-anomali", "Missar attacktyper den aldrig sett (reverse shell, path traversal).",
         "Riktig Wazuh-data + regelbaserat komplement för känd attacksyntax"),
        ("LLM-integration", "Ollama-latency 5–30 s — oacceptabelt för synkrona endpoints.",
         "Asynkron köhantering (BackgroundTasks) + caching av frekventa promptsvar"),
    ]

    for i, (title, limit, next_step) in enumerate(items):
        y = Inches(1.3 + i * 1.8)
        _box(slide, Inches(0.5), y, Inches(9), Inches(0.45),
             title, font_size=17, bold=True, color=ORANGE)
        _box(slide, Inches(0.5), y + Inches(0.45), Inches(9), Inches(0.55),
             f"⚠  {limit}", font_size=13, color=WHITE)
        _box(slide, Inches(0.5), y + Inches(1.0), Inches(9), Inches(0.55),
             f"→  Nästa steg: {next_step}", font_size=13, color=GREEN)

    _box(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.45),
         "Defense in depth: ClamAV + ML är starkare tillsammans än var för sig.",
         font_size=13, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)


def _conclusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _box(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
         "Slutsats", font_size=36, bold=True, color=TEAL)

    achievements = [
        ("✓", "3 tränade ML-modeller  ·  F1 0.875 på riktig CTI-data"),
        ("✓", "6 FastAPI-endpoints  ·  körs på Hetzner k3s i produktion"),
        ("✓", "104 automatiserade tester  ·  inkl. 8 outcome-tester för modellregression"),
        ("✓", "Adversarial-analys: poisoning, evasion och mimicry-attack"),
        ("✓", "Defense in depth: ClamAV + ML-klassificering + loggövervakning"),
    ]

    for i, (check, text) in enumerate(achievements):
        y = Inches(1.4 + i * 0.85)
        _box(slide, Inches(0.5), y, Inches(0.5), Inches(0.6),
             check, font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        _box(slide, Inches(1.1), y + Inches(0.08), Inches(8.4), Inches(0.55),
             text, font_size=15, color=WHITE)

    _box(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.7),
         '"Vi vet exakt var modellerna är svaga och varför.\nDet är mer värt i ett säkerhetssystem än ett F1 på 0.99 ni inte kan förklara."',
         font_size=15, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)


def main(out: Path = typer.Option(Path("docs/sentinel-ml-presentation.pptx"))) -> None:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _title_slide(prs)
    _problem_slide(prs)
    _architecture_slide(prs)
    _results_slide(prs)
    _adversarial_slide(prs)
    _limitations_slide(prs)
    _conclusion_slide(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    typer.echo(f"Presentation sparad: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    typer.run(main)
